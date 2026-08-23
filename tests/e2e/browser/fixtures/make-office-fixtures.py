#!/usr/bin/env python3
"""Regenerate the committed Office fixtures used by file-preview-kinds.spec.ts.

These are REAL OOXML binaries on purpose: a hand-built minimal zip is exactly
the "unsupported variant" the client-side renderers (docx-preview / SheetJS /
pptx-preview) reject, so a synthetic fixture would test the error path instead
of the preview.

    pip install python-docx openpyxl python-pptx
    python3 tests/e2e/browser/fixtures/make-office-fixtures.py

Keep the MARKER strings in sync with the spec's assertions. office-crafted.pptx
is deliberately INVALID (a literal '<' inside slide text, which no valid OOXML
writer emits) and exists only to prove a malformed deck cannot hang the tab.
"""
import pathlib
import zipfile

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

HERE = pathlib.Path(__file__).parent


def make_docx() -> None:
    d = Document()
    d.add_heading('Office preview fixture', level=1)
    d.add_paragraph('WALNUT DOCX MARKER')
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = 'Quarter'
    t.cell(0, 1).text = 'Revenue'
    t.cell(1, 0).text = 'Q1'
    t.cell(1, 1).text = '42'
    d.save(HERE / 'office-doc.docx')


def make_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = 'Data'
    ws['A1'] = 'WALNUT XLSX MARKER'
    ws['A2'], ws['B2'] = 'Item', 'Count'
    ws['A3'], ws['B3'] = 'apples', 7
    ws['A4'], ws['B4'] = 'pears', 3
    wb.create_sheet('Second')['A1'] = 'SECOND SHEET MARKER'
    wb.save(HERE / 'office-sheet.xlsx')


def make_pptx() -> None:
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[1])
    slide.shapes.title.text = 'WALNUT PPTX MARKER'
    slide.placeholders[1].text = 'Rendered by pptx-preview'
    p.save(HERE / 'office-slides.pptx')


def make_crafted_pptx() -> None:
    """Copy office-slides.pptx with a literal '<' injected into slide text."""
    src, dst = HERE / 'office-slides.pptx', HERE / 'office-crafted.pptx'
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == 'ppt/slides/slide1.xml':
                data = data.replace(b'WALNUT PPTX MARKER', b'<img src=x onerror=alert(1)>')
            zout.writestr(item, data)


if __name__ == '__main__':
    make_docx()
    make_xlsx()
    make_pptx()
    make_crafted_pptx()
    for f in sorted(HERE.glob('office-*')):
        print(f.name, f.stat().st_size)
